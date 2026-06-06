# -*- coding: utf-8 -*-
"""Build the UM iFEST 2026 FYP poster from the official template.

Preserves the master background (UM header band, footer 8-11 June 2026 band),
the top-right iFEST logo and the paper-ID textbox; removes the imported
Chinese-linguistics content and lays out a clean 2-column UEBA poster.
"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

SRC = r"C:\Users\yyzhe\OneDrive\Desktop\Learning Materials\FYP\UG FYP Poster Template 2026 with template.pptx"
# Fallback: when the OneDrive template is open in PowerPoint it is exclusively
# locked; use a previously saved local copy instead.
_SRC_FALLBACK = r"C:\Users\yyzhe\OneDrive\Documents\fyp-ueba\_template_local.pptx"
import zipfile as _zip
if not _zip.is_zipfile(SRC) and _zip.is_zipfile(_SRC_FALLBACK):
    SRC = _SRC_FALLBACK
OUT = r"C:\Users\yyzhe\OneDrive\Documents\fyp-ueba\Argus_UEBA_Poster.pptx"

PURPLE   = RGBColor(0x4B, 0x2E, 0x83)   # UM deep purple
PURPLE2  = RGBColor(0x6A, 0x3D, 0xB5)   # lighter accent
DARK     = RGBColor(0x22, 0x22, 0x28)
GREY     = RGBColor(0x55, 0x55, 0x5E)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHTBG  = RGBColor(0xF2, 0xEE, 0xF9)
ACCENT   = RGBColor(0xC8, 0x10, 0x6E)   # magenta highlight

prs = Presentation(SRC)
slide = prs.slides[0]

# ---- 1. strip the imported content, keep branding (object 50 id, 51 logo) ----
KEEP_IDS = {50, 51}
for sh in list(slide.shapes):
    if sh.shape_id not in KEEP_IDS:
        sh._element.getparent().remove(sh._element)

# update the paper-id strip text
for sh in slide.shapes:
    if sh.shape_id == 50 and sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.text = ""
        sh.text_frame.paragraphs[0].runs and None
        # rewrite cleanly
        tf = sh.text_frame
        tf.paragraphs[0].text = "UM iFEST 2026  ·  UG FYP"

# ---- helpers ----------------------------------------------------------------
def textbox(l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    return tb, tf

def setpara(p, text, size, color=DARK, bold=False, italic=False,
            align=PP_ALIGN.LEFT, space_after=4, bullet=False, lh=1.05):
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    p.line_spacing = lh
    r = p.add_run(); r.text = (("•  " if bullet else "") + text)
    f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = "Calibri"
    return p

def section_header(l, t, w, label, h=0.62):
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(l), Inches(t), Inches(w), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = PURPLE
    bar.line.fill.background()
    bar.shadow.inherit = False
    tf = bar.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    setpara(p, label, 23, WHITE, bold=True, align=PP_ALIGN.LEFT, space_after=0)
    p.runs[0].text = "  " + label
    return t + h + 0.12

def body_block(l, t, w, lines):
    """lines: list of dicts {text,size,bold,italic,color,align,bullet,space_after}"""
    # estimate height generously; textbox auto-grows visually
    tb, tf = textbox(l, t, w, 4.0)
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        setpara(p, ln["text"], ln.get("size", 15.5),
                color=ln.get("color", DARK), bold=ln.get("bold", False),
                italic=ln.get("italic", False), align=ln.get("align", PP_ALIGN.JUSTIFY),
                space_after=ln.get("space_after", 5), bullet=ln.get("bullet", False),
                lh=ln.get("lh", 1.04))
    return tb

def add_image(l, t, w, path):
    im = Image.open(path); ar = im.size[1] / im.size[0]
    h = w * ar
    slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    return t + h

def caption(l, t, w, text):
    tb, tf = textbox(l, t, w, 0.4)
    setpara(tf.paragraphs[0], text, 12.5, GREY, italic=True,
            align=PP_ALIGN.CENTER, space_after=0)
    return t + 0.38

# ============================================================================
PAGE_W = Emu(prs.slide_width).inches
M = 0.55
COL_W = 6.95
LX = M
RX = M + COL_W + 0.45   # 7.95

# ---- Title banner ----
tb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(M), Inches(1.12), Inches(PAGE_W - 2*M), Inches(2.15))
tb.fill.solid(); tb.fill.fore_color.rgb = PURPLE
tb.line.color.rgb = PURPLE2; tb.line.width = Pt(1.5)
tb.shadow.inherit = False
tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.4); tf.margin_right = Inches(0.4)
setpara(tf.paragraphs[0],
        "Cloud-Native UEBA Platform for Insider Threat Detection",
        38, WHITE, bold=True, align=PP_ALIGN.CENTER, space_after=4, lh=1.0)
p2 = tf.add_paragraph()
setpara(p2, "Unsupervised, explainable anomaly detection across four real cloud log sources",
        18, RGBColor(0xE3, 0xD8, 0xF5), italic=True, align=PP_ALIGN.CENTER, space_after=0, lh=1.0)

# ---- Authors / affiliation ----
_, atf = textbox(M, 3.40, PAGE_W - 2*M, 0.95, anchor=MSO_ANCHOR.TOP)
setpara(atf.paragraphs[0],
        "Yap Zhe Cheng  (23004982)      ·      Supervisor: Firdaus Sahran",
        20, DARK, bold=True, align=PP_ALIGN.CENTER, space_after=2)
pa = atf.add_paragraph()
setpara(pa, "Faculty of Computer Science and Information Technology, Universiti Malaya",
        15, GREY, italic=True, align=PP_ALIGN.CENTER, space_after=0)

# ============================ LEFT COLUMN ===================================
y = 4.55
y = section_header(LX, y, COL_W, "INTRODUCTION")
body_block(LX, y, COL_W, [
    {"text": "Insider threats drive ~60% of breaches and cost organisations an average of "
             "USD 17.4 million per year. Because insiders act with legitimate credentials, "
             "signature-based defences are structurally blind to them."},
    {"text": "Existing open-source UEBA tools (Wazuh, OpenUBA, HELK) are rule-based, "
             "single-source, and opaque — they cannot explain why a user is flagged. "
             "This project closes four gaps at once: multi-cloud normalisation, adaptive "
             "baselining, label-free detection, and per-alert explainability.",
     "space_after": 2},
])
y += 2.30

y = section_header(LX, y, COL_W, "OBJECTIVES")
body_block(LX, y, COL_W, [
    {"text": "Normalise heterogeneous multi-cloud logs into one unified 8-field schema.", "bullet": True},
    {"text": "Detect insider behaviour with unsupervised models — no labelled attacks.", "bullet": True},
    {"text": "Explain every alert via SHAP attribution and an LLM analyst assistant.", "bullet": True},
    {"text": "Evaluate against the CERT r4.2 insider-threat benchmark and real cloud data.", "bullet": True, "space_after": 2},
])
y += 1.85

y = section_header(LX, y, COL_W, "METHODOLOGY  —  SYSTEM ARCHITECTURE")
body_block(LX, y, COL_W, [
    {"text": "A decoupled, layered pipeline: source-specific parsers feed one schema; per-user "
             "behavioural features and peer-group baselines feed two Autoencoders; SHAP and an "
             "LLM explain each score on a React dashboard, deployed via Docker + HuggingFace.",
     "size": 14.5, "space_after": 4},
])
y += 0.95
_dw = 6.7
y_img = add_image(LX + (COL_W - _dw) / 2, y, _dw, "results/report_png/architecture_logos.png")
caption(LX, y_img + 0.04, COL_W,
        "Fig. 1.  Layered pipeline: ingestion → features → detection → explainability → dashboard.")

# ============================ RIGHT COLUMN ==================================
y = 4.55
y = section_header(RX, y, COL_W, "RESULTS  &  DISCUSSION")

# results table
rows = [
    ["Model", "AUROC", "AUPRC", "F1"],
    ["Autoencoder (final)", "0.976", "0.851", "0.787"],
    ["Weighted IF+AE", "0.951", "0.740", "0.726"],
    ["Isolation Forest", "0.860", "0.206", "0.379"],
    ["Rule-based", "0.859", "0.207", "0.367"],
]
tbl_h = 1.9
gtbl = slide.shapes.add_table(len(rows), 4, Inches(RX), Inches(y),
                              Inches(COL_W), Inches(tbl_h)).table
gtbl.columns[0].width = Inches(3.05)
for c in range(1, 4):
    gtbl.columns[c].width = Inches((COL_W - 3.05) / 3)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = gtbl.cell(ri, ci)
        cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.05)
        cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        if ri == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = PURPLE
        elif ri == 1:
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHTBG
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        r = p.add_run(); r.text = val
        r.font.size = Pt(14.5)
        r.font.bold = (ri == 0) or (ri == 1)
        r.font.color.rgb = WHITE if ri == 0 else DARK
        r.font.name = "Calibri"
y += tbl_h + 0.08
caption(RX, y, COL_W, "Table 1.  Detector performance on CERT r4.2 (1,000 users, 70 insiders).")
y += 0.40

roc_w = 5.2
roc_l = RX + (COL_W - roc_w) / 2
y_img = add_image(roc_l, y, roc_w, "results/report_png/roc.png")
y = caption(RX, y_img + 0.02, COL_W,
            "Fig. 2.  ROC curves — Autoencoder vs. baselines.")
body_block(RX, y, COL_W, [
    {"text": "The label-free Autoencoder reaches ~90% of supervised accuracy (AUROC 0.976) and "
             "decisively beats rule and isolation baselines. A live-replay harness confirms the "
             "deployed /ingest path preserves this ranking (AUROC 0.917); on real flaws.cloud AWS "
             "data the cloud model reaches AUROC 0.724 against privilege-escalation attacks.",
     "size": 14.5, "space_after": 2},
])
y += 1.95

y = section_header(RX, y, COL_W, "CONCLUSIONS")
body_block(RX, y, COL_W, [
    {"text": "An unsupervised Autoencoder detects insiders at near-supervised accuracy with no labels.", "bullet": True},
    {"text": "Four real cloud sources unified; SHAP + LLM make every alert explainable.", "bullet": True},
    {"text": "Deployed live (Docker + HuggingFace Space + Vercel).", "bullet": True},
    {"text": "Future: Transformer sequence model, MITRE ATT&CK mapping, streaming ingestion.", "bullet": True, "space_after": 2},
])
y += 1.95

y = section_header(RX, y, COL_W, "REFERENCES")
body_block(RX, y, COL_W, [
    {"text": "Inayat et al. (2024). Insider threat mitigation: systematic review. Ain Shams Eng. J.", "size": 12.5, "space_after": 2, "lh": 1.0},
    {"text": "Kotb et al. (2025). Deep synthesis-based insider intrusion detection. Sci. Reports.", "size": 12.5, "space_after": 2, "lh": 1.0},
    {"text": "Datta et al. (2021). Real-time threat detection in UEBA. IEEE IEMENTech.", "size": 12.5, "space_after": 2, "lh": 1.0},
    {"text": "Verizon (2025). Data Breach Investigations Report.", "size": 12.5, "space_after": 0, "lh": 1.0},
])

# ============================ FULL-WIDTH VALUE BAND =========================
# Highlights the project's core contributions / novelty / impact at a glance.
FULL_W = PAGE_W - 2 * M
by = 18.92
by = section_header(M, by, FULL_W, "KEY CONTRIBUTIONS  &  IMPACT", h=0.48)

cards = [
    ("AUROC 0.976",      "Label-free detection — near-supervised accuracy with zero labelled attacks."),
    ("4 cloud sources",  "AWS · Azure · Cloudflare · GitHub unified into one 8-field schema."),
    ("Explainable",      "SHAP attribution + LLM analyst give the reason behind every alert."),
    ("Live AUROC 0.917", "Production /ingest validated; deployed on Docker + HF Space + Vercel."),
]
gap = 0.18
cw = (FULL_W - gap * (len(cards) - 1)) / len(cards)
for i, (metric, desc) in enumerate(cards):
    cx = M + i * (cw + gap)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(cx), Inches(by), Inches(cw), Inches(0.96))
    card.fill.solid(); card.fill.fore_color.rgb = LIGHTBG
    card.line.color.rgb = PURPLE2; card.line.width = Pt(1.25)
    card.shadow.inherit = False
    tf = card.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    setpara(tf.paragraphs[0], metric, 19, ACCENT, bold=True,
            align=PP_ALIGN.CENTER, space_after=2, lh=1.0)
    pd = tf.add_paragraph()
    setpara(pd, desc, 12.5, DARK, align=PP_ALIGN.CENTER, space_after=0, lh=1.0)

prs.save(OUT)
print("WROTE", OUT)
